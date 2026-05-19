"""Shared runtime for boundary-first SkillGraph artifacts."""

from __future__ import annotations

import ast
import hashlib
import json
import os
import re
import shlex
import subprocess
import sys
import zipfile
from pathlib import Path
from typing import Any

_SOLVER_BOUNDARIES = {"workflow_solver", "jit_solver"}
_GENERIC_EXECUTABLE_OPERATOR_KINDS = {
    "docx_inspect",
    "docx_replace_text",
    "docx_marked_blocks",
    "docx_validate_text",
    "pptx_inventory",
    "pptx_update_text_boxes",
    "pptx_add_reference_slide",
    "pptx_validate_file",
}


def run_boundary_artifact(
    spec: dict[str, Any],
    state: dict[str, Any],
) -> dict[str, Any]:
    """Execute a boundary artifact and return the canonical runtime envelope."""

    inputs = _inputs(state)
    boundary = str(spec.get("boundary") or "adapter")
    planlet = _select_planlet(spec, inputs)
    if planlet is not None and (inputs.get("planlet") or not inputs.get("operation")):
        missing = _missing_planlet_bindings(planlet, inputs)
        if missing:
            return _completed(
                spec,
                runtime_status="planlet_guidance",
                contribution_type="planlet_guidance",
                final_output=_planlet_guidance(planlet, missing),
                selected_operator=None,
                selected_planlet=planlet,
                requires_agent_continuation=True,
            )
        return _execute_planlet(spec, planlet, inputs)

    if boundary == "guidance":
        return _guidance_result(spec, inputs)

    operator = _select_operator(spec, inputs)
    if operator is None:
        return _blocked(
            spec,
            "no_match",
            "No boundary operator matched the request.",
            {"available_operations": _operator_names(spec)},
        )

    if _network_blocked(spec, operator, inputs):
        return _blocked(
            spec,
            "network_risk",
            "Operation requires network access and boundary policy blocks network by default.",
            {"operator": operator.get("name") or operator.get("id")},
        )

    if _is_lossless_reference_operator(operator):
        execution = _execute_lossless_reference(spec, operator, inputs)
        if execution["status"] != "completed":
            return _blocked(
                spec,
                execution.get("deopt_reason") or "reference_lookup_failed",
                execution.get("error") or "Lossless reference lookup failed.",
                execution,
            )
        return _completed(
            spec,
            runtime_status="guidance_only",
            contribution_type="retrieval_guidance",
            final_output=execution.get("final_output"),
            selected_operator=operator,
            requires_agent_continuation=True,
            debug_trace={"execution": execution},
        )

    if boundary == "adapter" and not _truthy(inputs.get("execute")):
        return _completed(
            spec,
            runtime_status="guidance_only",
            contribution_type="retrieval_guidance",
            final_output=_operator_guidance(operator),
            selected_operator=operator,
            requires_agent_continuation=True,
        )

    if operator.get("kind") == "reference_section":
        return _completed(
            spec,
            runtime_status="guidance_only",
            contribution_type="retrieval_guidance",
            final_output=_operator_guidance(operator),
            selected_operator=operator,
            requires_agent_continuation=True,
        )

    validation = _validate_operator_arguments(operator, inputs)
    if validation["errors"]:
        return _blocked(
            spec,
            "invalid_arguments",
            f"Invalid arguments for operation {operator.get('name') or operator.get('id')}.",
            validation,
        )

    execution = _execute_operator(spec, operator, inputs)
    if execution["status"] != "completed":
        return _blocked(
            spec,
            execution.get("deopt_reason") or "execution_failed",
            execution.get("error") or "Boundary operator execution failed.",
            execution,
        )

    requires_continuation = boundary not in _SOLVER_BOUNDARIES
    return _completed(
        spec,
        runtime_status="executed",
        contribution_type="execution" if not requires_continuation else "typed_evidence",
        final_output=execution.get("final_output"),
        selected_operator=operator,
        requires_agent_continuation=requires_continuation,
        step_results={
            str(operator.get("id") or operator.get("name")): {
                "executor": execution.get("executor"),
                "returncode": execution.get("returncode"),
                "stdout": execution.get("stdout"),
                "stderr": execution.get("stderr"),
                "output": execution.get("final_output"),
            }
        },
        debug_trace={"execution": execution},
    )


def normalize_boundary_envelope(
    result: dict[str, Any],
    spec: dict[str, Any] | None,
) -> dict[str, Any]:
    """Backfill boundary runtime fields onto thin or legacy generated results."""

    if not isinstance(result, dict):
        return _blocked(spec or {}, "invalid_result", f"Artifact returned {type(result).__name__}.")
    boundary_spec = spec or {}
    boundary = str(boundary_spec.get("boundary") or result.get("boundary") or "adapter")
    enriched = dict(result)
    enriched.setdefault("boundary", boundary)
    enriched.setdefault("artifact_boundary", _summary(boundary_spec))
    enriched["can_produce_final_answer"] = boundary in _SOLVER_BOUNDARIES
    if boundary not in _SOLVER_BOUNDARIES:
        enriched["requires_agent_continuation"] = True
        if enriched.get("status") == "completed" and enriched.get("runtime_status") == "executed":
            enriched.setdefault("contribution_type", "typed_evidence")
    else:
        enriched.setdefault("requires_agent_continuation", enriched.get("status") != "completed")
    enriched.setdefault(
        "runtime_status",
        "executed" if enriched.get("status") == "completed" else "blocked",
    )
    enriched.setdefault("errors", [])
    enriched.setdefault("produced_outputs", enriched.get("outputs", []))
    return enriched


def _execute_operator(
    spec: dict[str, Any],
    operator: dict[str, Any],
    inputs: dict[str, Any],
) -> dict[str, Any]:
    kind = str(operator.get("kind") or "")
    try:
        if kind in _GENERIC_EXECUTABLE_OPERATOR_KINDS:
            return _execute_generic_operator(kind, operator, inputs)
        if kind == "script_asset" and operator.get("script_path"):
            return _execute_script(spec, operator, inputs)
        if kind == "shell_command" and operator.get("commands"):
            return _execute_shell(spec, operator, inputs)
    except ValueError as exc:
        return {
            "status": "blocked",
            "deopt_reason": "invalid_arguments",
            "error": str(exc),
        }
    except Exception as exc:
        return {
            "status": "blocked",
            "deopt_reason": "execution_exception",
            "error": f"{type(exc).__name__}: {exc}",
        }
    return {
        "status": "blocked",
        "deopt_reason": "unsupported_operator",
        "error": f"Unsupported boundary operator kind: {kind or 'unknown'}",
    }


def _execute_generic_operator(
    kind: str,
    operator: dict[str, Any],
    inputs: dict[str, Any],
) -> dict[str, Any]:
    dispatch = {
        "docx_inspect": _docx_inspect,
        "docx_replace_text": _docx_replace_text,
        "docx_marked_blocks": _docx_marked_blocks,
        "docx_validate_text": _docx_validate_text,
        "pptx_inventory": _pptx_inventory,
        "pptx_update_text_boxes": _pptx_update_text_boxes,
        "pptx_add_reference_slide": _pptx_add_reference_slide,
        "pptx_validate_file": _pptx_validate_file,
    }
    output = dispatch[kind](operator, inputs)
    return _completed_execution(output, executor=kind)


def _is_lossless_reference_operator(operator: dict[str, Any]) -> bool:
    return str(operator.get("kind") or "") in {
        "skill_asset_list",
        "skill_asset_get",
        "skill_doc_search",
    }


def _execute_lossless_reference(
    spec: dict[str, Any],
    operator: dict[str, Any],
    inputs: dict[str, Any],
) -> dict[str, Any]:
    kind = str(operator.get("kind") or "")
    try:
        if kind == "skill_asset_list":
            output = _lossless_list_assets(spec, inputs)
        elif kind == "skill_asset_get":
            output = _lossless_get_asset(spec, inputs)
        elif kind == "skill_doc_search":
            output = _lossless_search_docs(spec, inputs)
        else:
            return {
                "status": "blocked",
                "deopt_reason": "unsupported_reference_operator",
                "error": f"Unsupported lossless reference operator: {kind}",
            }
    except ValueError as exc:
        return {
            "status": "blocked",
            "deopt_reason": "invalid_reference_request",
            "error": str(exc),
            "executor": kind,
        }
    except Exception as exc:
        return {
            "status": "blocked",
            "deopt_reason": "reference_exception",
            "error": f"{type(exc).__name__}: {exc}",
            "executor": kind,
        }
    return _completed_execution(output, executor=kind)


def _lossless_list_assets(spec: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    args = _arguments(inputs)
    asset_type = str(args.get("asset_type") or "").strip().lower()
    assets = _lossless_assets(spec)
    if asset_type:
        assets = [
            asset
            for asset in assets
            if str(asset.get("asset_type") or "").strip().lower() == asset_type
        ]
    return {
        "assets": assets,
        "asset_count": len(assets),
        "lossless_reference": _lossless_summary(spec),
    }


def _lossless_get_asset(spec: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    args = _arguments(inputs)
    path_value = str(args.get("path") or args.get("file") or "").strip()
    if not path_value:
        raise ValueError("Missing required argument: path")
    asset = _find_lossless_asset(spec, path_value)
    skill_root = _lossless_skill_root(inputs)
    asset_path = _resolve_skill_asset_path(skill_root, path_value)
    if not asset_path.exists():
        raise ValueError(f"Skill asset not found in sandbox: {path_value}")
    if not bool(asset.get("is_text")):
        return {
            "path": asset["path"],
            "asset_type": asset.get("asset_type", ""),
            "is_text": False,
            "size_bytes": asset_path.stat().st_size,
            "sha256": _sha256_file(asset_path),
            "sandbox_path": str(asset_path),
            "lossless_reference": _lossless_summary(spec),
        }
    text = asset_path.read_text(encoding="utf-8")
    lines = text.splitlines()
    start_line = max(1, _safe_int(args.get("start_line"), default=1))
    end_line = min(len(lines), max(start_line, _safe_int(args.get("end_line"), default=len(lines))))
    max_chars = _safe_int(args.get("max_chars"), default=0)
    selected = "\n".join(lines[start_line - 1 : end_line]) if lines else ""
    truncated = False
    if max_chars > 0 and len(selected) > max_chars:
        selected = selected[:max_chars]
        truncated = True
    return {
        "path": asset["path"],
        "asset_type": asset.get("asset_type", ""),
        "is_text": True,
        "size_bytes": asset_path.stat().st_size,
        "sha256": _sha256_file(asset_path),
        "line_count": len(lines),
        "start_line": start_line,
        "end_line": end_line,
        "content": selected,
        "truncated": truncated,
        "lossless_reference": _lossless_summary(spec),
    }


def _lossless_search_docs(spec: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    args = _arguments(inputs)
    query = str(
        args.get("query")
        or inputs.get("request")
        or inputs.get("task")
        or inputs.get("topic")
        or ""
    ).strip()
    if not query:
        raise ValueError("Missing required argument: query")
    max_results = max(1, _safe_int(args.get("max_results"), default=8))
    terms = [term.lower() for term in re.findall(r"[A-Za-z0-9_./-]{2,}", query)]
    if not terms:
        terms = [query.lower()]
    skill_root = _lossless_skill_root(inputs)
    results: list[dict[str, Any]] = []
    for asset in _lossless_assets(spec):
        if not bool(asset.get("is_text")):
            continue
        path = str(asset.get("path") or "")
        asset_path = _resolve_skill_asset_path(skill_root, path)
        if not asset_path.exists():
            continue
        try:
            lines = asset_path.read_text(encoding="utf-8").splitlines()
        except UnicodeDecodeError:
            continue
        for idx, line in enumerate(lines, start=1):
            lower = line.lower()
            score = sum(1 for term in terms if term in lower)
            if score <= 0:
                continue
            results.append(
                {
                    "path": path,
                    "line": idx,
                    "score": score,
                    "snippet": line.strip()[:500],
                    "asset_type": asset.get("asset_type", ""),
                }
            )
    results.sort(
        key=lambda item: (
            -int(item.get("score") or 0),
            str(item.get("path") or ""),
            int(item.get("line") or 0),
        )
    )
    return {
        "query": query,
        "results": results[:max_results],
        "result_count": len(results[:max_results]),
        "total_matches": len(results),
        "lossless_reference": _lossless_summary(spec),
    }


def _lossless_reference(spec: dict[str, Any]) -> dict[str, Any]:
    scope = spec.get("scope") if isinstance(spec, dict) else {}
    lossless = scope.get("lossless_reference") if isinstance(scope, dict) else {}
    return lossless if isinstance(lossless, dict) else {}


def _lossless_summary(spec: dict[str, Any]) -> dict[str, Any]:
    lossless = _lossless_reference(spec)
    return {
        "strict_validated": bool(lossless.get("strict_validated")),
        "package_hash": str(lossless.get("package_hash") or ""),
        "capsule_hash": str(lossless.get("capsule_hash") or ""),
        "asset_count": int(lossless.get("asset_count") or 0),
        "text_asset_count": int(lossless.get("text_asset_count") or 0),
        "binary_asset_count": int(lossless.get("binary_asset_count") or 0),
    }


def _lossless_assets(spec: dict[str, Any]) -> list[dict[str, Any]]:
    assets = _lossless_reference(spec).get("assets") or []
    return [asset for asset in assets if isinstance(asset, dict) and asset.get("path")]


def _find_lossless_asset(spec: dict[str, Any], path_value: str) -> dict[str, Any]:
    normalized = _normalize_skill_asset_path(path_value)
    for asset in _lossless_assets(spec):
        if _normalize_skill_asset_path(str(asset.get("path") or "")) == normalized:
            return asset
    raise ValueError(f"Unknown skill asset: {path_value}")


def _lossless_skill_root(inputs: dict[str, Any]) -> Path:
    raw = str(inputs.get("sandbox_dir") or "").strip()
    if raw:
        skill_root = Path(raw) / "skill"
        if skill_root.exists():
            return skill_root
    raise ValueError("Lossless reference requires sandbox_dir with a skill copy.")


def _resolve_skill_asset_path(skill_root: Path, path_value: str) -> Path:
    normalized = _normalize_skill_asset_path(path_value)
    candidate = (skill_root / normalized).resolve()
    root = skill_root.resolve()
    if candidate != root and root not in candidate.parents:
        raise ValueError(f"Skill asset path escapes sandbox: {path_value}")
    return candidate


def _normalize_skill_asset_path(path_value: str) -> str:
    path = path_value.strip().replace("\\", "/").lstrip("/")
    parts = [part for part in path.split("/") if part not in {"", "."}]
    if any(part == ".." for part in parts):
        raise ValueError(f"Invalid skill asset path: {path_value}")
    return "/".join(parts)


def _completed_execution(output: Any, *, executor: str) -> dict[str, Any]:
    return {
        "status": "completed",
        "final_output": output,
        "returncode": 0,
        "stdout": json.dumps(output, ensure_ascii=False),
        "stderr": "",
        "executor": executor,
    }


def _docx_inspect(operator: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    from docx import Document

    args = _arguments(inputs, operator)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    path = _required_workspace_path_alias(
        project_dir, args, "path", "docx_path", "input_path", "file"
    )
    doc = Document(str(path))
    paragraphs = [
        {"index": idx, "text": paragraph.text}
        for idx, paragraph in enumerate(doc.paragraphs)
        if paragraph.text
    ]
    tables = list(_docx_table_text_items(doc.tables))
    header_footer = []
    for section_idx, section in enumerate(doc.sections):
        for location, part in (("header", section.header), ("footer", section.footer)):
            for paragraph_idx, paragraph in enumerate(part.paragraphs):
                if paragraph.text:
                    header_footer.append(
                        {
                            "section": section_idx,
                            "location": location,
                            "paragraph": paragraph_idx,
                            "text": paragraph.text,
                        }
                    )
    all_text = "\n".join(
        [item["text"] for item in paragraphs]
        + [item["text"] for item in tables]
        + [item["text"] for item in header_footer]
    )
    placeholders = sorted(set(re.findall(r"\{\{[^{}]{1,120}\}\}|<<[^<>]{1,120}>>", all_text)))
    return {
        "path": str(path),
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
        "paragraphs": paragraphs[:120],
        "tables": tables[:120],
        "headers_footers": header_footer[:80],
        "placeholders": placeholders,
    }


def _docx_replace_text(operator: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    from docx import Document

    args = _arguments(inputs, operator)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    input_path = _required_workspace_path(project_dir, args, "input_path")
    output_path = _required_workspace_path(project_dir, args, "output_path")
    replacements = _mapping_arg(args, "replacements")
    include_headers_footers = _truthy_arg(args, "include_headers_footers", default=True)
    include_tables = _truthy_arg(args, "include_tables", default=True)
    doc = Document(str(input_path))
    changed = 0
    for paragraph in _docx_paragraphs(
        doc,
        include_tables=include_tables,
        include_headers_footers=include_headers_footers,
    ):
        changed += _replace_paragraph_text(paragraph, replacements)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    return {
        "input_path": str(input_path),
        "output_path": str(output_path),
        "replacement_count": len(replacements),
        "changed_paragraphs": changed,
        "exists": output_path.exists(),
    }


def _docx_marked_blocks(operator: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    from docx import Document

    args = _arguments(inputs, operator)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    input_path = _required_workspace_path(project_dir, args, "input_path")
    output_path = _required_workspace_path(project_dir, args, "output_path")
    markers = _marker_specs(args)
    conditions = _mapping_arg(args, "conditions", required=False)
    doc = Document(str(input_path))
    changed = 0
    for marker in markers:
        name = str(marker.get("name") or marker.get("id") or marker.get("start") or "")
        start = str(marker.get("start") or marker.get("start_marker") or "")
        end = str(marker.get("end") or marker.get("end_marker") or "")
        if not start or not end:
            raise ValueError("Each marker requires start and end")
        keep = marker.get("keep")
        if keep is None and name:
            keep = conditions.get(name)
        keep_block = bool(keep)
        in_block = False
        for paragraph in doc.paragraphs:
            text = paragraph.text
            if start in text:
                in_block = True
                changed += 1
                if keep_block:
                    paragraph.text = text.replace(start, "").replace(end, "")
                    if end in text:
                        in_block = False
                else:
                    paragraph.text = ""
                    if end in text:
                        in_block = False
                continue
            if in_block:
                changed += 1
                paragraph.text = text.replace(end, "") if keep_block else ""
                if end in text:
                    in_block = False
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    return {
        "input_path": str(input_path),
        "output_path": str(output_path),
        "markers_processed": len(markers),
        "changed_paragraphs": changed,
        "exists": output_path.exists(),
    }


def _docx_validate_text(operator: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    from docx import Document

    args = _arguments(inputs, operator)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    path = _required_workspace_path_alias(
        project_dir, args, "path", "docx_path", "input_path", "file"
    )
    doc = Document(str(path))
    text = "\n".join(
        paragraph.text
        for paragraph in _docx_paragraphs(
            doc,
            include_tables=True,
            include_headers_footers=True,
        )
        if paragraph.text
    )
    must_contain = _list_arg(args, "must_contain")
    must_not_contain = _list_arg(args, "must_not_contain")
    missing = [item for item in must_contain if item not in text]
    forbidden_found = [item for item in must_not_contain if item in text]
    return {
        "path": str(path),
        "ok": not missing and not forbidden_found,
        "missing": missing,
        "forbidden_found": forbidden_found,
        "character_count": len(text),
    }


def _pptx_inventory(operator: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation

    args = _arguments(inputs, operator)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    input_pptx = _required_workspace_path_alias(
        project_dir, args, "input_pptx", "pptx_path", "path", "input", "file"
    )
    output_json = _optional_workspace_path(project_dir, args, "output_json")
    include_positions = _truthy_arg(args, "include_positions", default=True)
    prs = Presentation(str(input_pptx))
    slides: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            item = {
                "shape_id": int(shape.shape_id),
                "name": str(shape.name),
                "has_text": bool(getattr(shape, "has_text_frame", False)),
                "text": shape.text if getattr(shape, "has_text_frame", False) else "",
            }
            if include_positions:
                item.update(
                    {
                        "left": int(shape.left),
                        "top": int(shape.top),
                        "width": int(shape.width),
                        "height": int(shape.height),
                    }
                )
            shapes.append(item)
        slides.append({"slide_index": slide_idx, "shape_count": len(shapes), "shapes": shapes})
    output = {"input_pptx": str(input_pptx), "slide_count": len(slides), "slides": slides}
    if output_json:
        output_json.parent.mkdir(parents=True, exist_ok=True)
        output_json.write_text(json.dumps(output, indent=2, ensure_ascii=False), encoding="utf-8")
        output["output_json"] = str(output_json)
    return output


def _pptx_update_text_boxes(operator: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    args = _arguments(inputs, operator)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    input_pptx = _required_workspace_path_alias(
        project_dir, args, "input_pptx", "pptx_path", "input_path", "path", "input", "file"
    )
    output_pptx = _required_workspace_path_alias(
        project_dir, args, "output_pptx", "output_path", "output", "destination", "dest"
    )
    if "edits" not in args and "updates" in args:
        args["edits"] = args["updates"]
    if "edits" not in args and "text_boxes" in args:
        args["edits"] = args["text_boxes"]
    if "edits" not in args and "targets" in args:
        targets = _list_arg(args, "targets", required=True)
        shared = {
            key: value
            for key, value in args.items()
            if key
            not in {
                "input_pptx",
                "pptx_path",
                "input_path",
                "path",
                "input",
                "file",
                "output_pptx",
                "output_path",
                "output",
                "destination",
                "dest",
                "targets",
                "text_boxes",
                "edits",
                "updates",
            }
        }
        args["edits"] = [
            {**target, **shared} if isinstance(target, dict) else target for target in targets
        ]
    edits = _list_arg(args, "edits", required=True)
    prs = Presentation(str(input_pptx))
    applied = 0
    for edit in edits:
        if not isinstance(edit, dict):
            continue
        nested_style = edit.get("style")
        if isinstance(nested_style, dict):
            edit = {**edit, **nested_style}
        nested_format = edit.get("format")
        if isinstance(nested_format, dict):
            edit = {**edit, **nested_format}
        for slide_idx, slide in enumerate(prs.slides, start=1):
            if edit.get("slide_index") and int(edit["slide_index"]) != slide_idx:
                continue
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if not _pptx_shape_matches(shape, edit):
                    continue
                if "new_text" in edit:
                    shape.text = str(edit["new_text"])
                elif "text" in edit and not any(
                    key in edit for key in ("shape_id", "contains", "match_text")
                ):
                    shape.text = str(edit["text"])
                for attr in ("left", "top", "width", "height"):
                    if attr in edit:
                        setattr(shape, attr, int(edit[attr]))
                _apply_pptx_positioning(shape, prs, edit)
                _apply_pptx_text_frame_options(shape, edit)
                if "font_size" in edit:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(float(edit["font_size"]))
                if any(
                    key in edit for key in ("font_name", "font_color", "bold", "alignment", "align")
                ):
                    _apply_pptx_font_style(shape, edit, RGBColor, Pt, PP_ALIGN)
                applied += 1
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    return {
        "input_pptx": str(input_pptx),
        "output_pptx": str(output_pptx),
        "requested_edits": len(edits),
        "applied_edits": applied,
        "exists": output_pptx.exists(),
    }


def _pptx_add_reference_slide(operator: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    args = _arguments(inputs, operator)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    input_pptx = _required_workspace_path_alias(
        project_dir, args, "input_pptx", "pptx_path", "input_path", "path", "input", "file"
    )
    output_pptx = _required_workspace_path_alias(
        project_dir, args, "output_pptx", "output_path", "output", "destination", "dest"
    )
    if "items" not in args and "references" in args:
        args["items"] = args["references"]
    if "items" not in args and "titles" in args:
        args["items"] = args["titles"]
    items = [str(item) for item in _list_arg(args, "items", required=True)]
    if _truthy_arg(args, "deduplicate", default=True):
        deduped: list[str] = []
        seen: set[str] = set()
        for item in items:
            key = " ".join(item.split()).strip().lower()
            if key and key not in seen:
                seen.add(key)
                deduped.append(item)
        items = deduped
    title = str(args.get("title") or "References")
    numbered = _truthy_arg(args, "numbered", default=True)
    prs = Presentation(str(input_pptx))
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9.0), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].runs[0].font.size = Pt(28)
    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.05), Inches(8.8), Inches(5.8))
    frame = body.text_frame
    frame.clear()
    for idx, item in enumerate(items, start=1):
        paragraph = frame.paragraphs[0] if idx == 1 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(14)
        if numbered:
            _set_pptx_auto_number(paragraph)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    return {
        "input_pptx": str(input_pptx),
        "output_pptx": str(output_pptx),
        "slide_count": len(prs.slides),
        "items_added": len(items),
        "exists": output_pptx.exists(),
    }


def _pptx_validate_file(operator: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation

    args = _arguments(inputs, operator)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    path = _required_workspace_path_alias(
        project_dir, args, "path", "pptx_path", "input_pptx", "file"
    )
    is_zip = zipfile.is_zipfile(path)
    prs = Presentation(str(path))
    return {
        "path": str(path),
        "ok": bool(is_zip),
        "is_zip": bool(is_zip),
        "slide_count": len(prs.slides),
    }


def _execute_script(
    spec: dict[str, Any],
    operator: dict[str, Any],
    inputs: dict[str, Any],
) -> dict[str, Any]:
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    sandbox_dir = Path(str(inputs.get("sandbox_dir") or ""))
    script_path = str(operator.get("script_path") or "")
    candidate_paths = [
        sandbox_dir / "skill" / script_path,
        project_dir / script_path,
        Path(script_path),
    ]
    script = next((path for path in candidate_paths if path.exists()), None)
    if script is None:
        return {
            "status": "blocked",
            "deopt_reason": "missing_script",
            "error": f"Script not found: {script_path}",
        }
    args = _arguments(inputs, operator)
    argv = [_python_executable(project_dir), str(script)]
    bindings = dict(operator.get("input_bindings") or {})
    input_fields = operator.get("inputs") or []
    if not input_fields:
        positional = _script_positional_from_context(args, inputs)
        if positional:
            argv.append(positional)
    for field in input_fields:
        name = str(field.get("name") or "")
        if not name or name not in args:
            if field.get("required"):
                return {
                    "status": "blocked",
                    "deopt_reason": "missing_argument",
                    "error": f"Missing required argument: {name}",
                }
            continue
        flag = str(bindings.get(name) or "")
        value = args[name]
        if isinstance(value, bool):
            if value and flag:
                argv.append(flag)
            elif value and not flag:
                argv.append(str(value))
            continue
        if flag:
            argv.extend([flag, str(value)])
        else:
            argv.append(str(value))
    return _run(argv, project_dir, _timeout(spec), executor="script_asset")


def _script_positional_from_context(args: dict[str, Any], inputs: dict[str, Any]) -> str:
    for key in (
        "file",
        "path",
        "input",
        "input_file",
        "input_path",
        "source",
        "source_file",
        "source_path",
        "data_file",
        "csv_file",
        "json_file",
        "bibtex_file",
        "stl_path",
    ):
        value = args.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    files = inputs.get("files")
    if isinstance(files, list):
        for value in files:
            if isinstance(value, str) and value.strip():
                return value.strip()
    return ""


def _execute_shell(
    spec: dict[str, Any],
    operator: dict[str, Any],
    inputs: dict[str, Any],
) -> dict[str, Any]:
    commands = [str(item) for item in operator.get("commands") or [] if item]
    if not commands:
        return {
            "status": "blocked",
            "deopt_reason": "missing_command",
            "error": "Operator has no shell command.",
        }
    args = _arguments(inputs, operator)
    command = _bind_command(commands[0], args)
    project_dir = Path(str(inputs.get("project_dir") or os.getcwd()))
    try:
        argv = shlex.split(command)
    except ValueError as exc:
        return {
            "status": "blocked",
            "deopt_reason": "invalid_command",
            "error": str(exc),
        }
    return _run(argv, project_dir, _timeout(spec), executor="shell_command")


def _run(
    argv: list[str],
    cwd: Path,
    timeout: int,
    *,
    executor: str,
) -> dict[str, Any]:
    try:
        completed = subprocess.run(
            argv,
            cwd=str(cwd),
            text=True,
            capture_output=True,
            timeout=timeout,
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        return {
            "status": "blocked",
            "deopt_reason": "timeout",
            "error": f"Execution exceeded {timeout}s.",
            "stdout": exc.stdout or "",
            "stderr": exc.stderr or "",
            "executor": executor,
        }
    except Exception as exc:
        return {
            "status": "blocked",
            "deopt_reason": "execution_exception",
            "error": f"{type(exc).__name__}: {exc}",
            "executor": executor,
        }
    if completed.returncode != 0:
        return {
            "status": "blocked",
            "deopt_reason": "nonzero_exit",
            "error": completed.stderr or completed.stdout or f"exit {completed.returncode}",
            "returncode": completed.returncode,
            "stdout": completed.stdout,
            "stderr": completed.stderr,
            "executor": executor,
        }
    return {
        "status": "completed",
        "final_output": _parse_output(completed.stdout),
        "returncode": completed.returncode,
        "stdout": completed.stdout,
        "stderr": completed.stderr,
        "executor": executor,
    }


def _guidance_result(spec: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any]:
    operator = _select_operator(spec, inputs)
    if operator is None:
        return _blocked(
            spec,
            "no_match",
            "No guidance section matched the request.",
            {"available_operations": _operator_names(spec)},
        )
    if _is_lossless_reference_operator(operator):
        execution = _execute_lossless_reference(spec, operator, inputs)
        if execution["status"] != "completed":
            return _blocked(
                spec,
                execution.get("deopt_reason") or "reference_lookup_failed",
                execution.get("error") or "Lossless reference lookup failed.",
                execution,
            )
        return _completed(
            spec,
            runtime_status="guidance_only",
            contribution_type="retrieval_guidance",
            final_output=execution.get("final_output"),
            selected_operator=operator,
            requires_agent_continuation=True,
            debug_trace={"execution": execution},
        )
    return _completed(
        spec,
        runtime_status="guidance_only",
        contribution_type="retrieval_guidance",
        final_output=_operator_guidance(operator),
        selected_operator=operator,
        requires_agent_continuation=True,
    )


def _select_planlet(spec: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any] | None:
    planlets = [item for item in spec.get("planlets") or [] if isinstance(item, dict)]
    if not planlets:
        return None
    requested = str(inputs.get("planlet") or "").strip().lower()
    if requested:
        for planlet in planlets:
            names = {
                str(planlet.get("id") or "").lower(),
                str(planlet.get("intent") or "").lower(),
            }
            if requested in names:
                return planlet
        for planlet in planlets:
            if requested in str(planlet.get("id") or "").lower():
                return planlet
    request_text = " ".join(
        str(value)
        for value in [inputs.get("request"), inputs.get("task"), inputs.get("topic")]
        if value
    ).lower()
    file_text = " ".join(str(value).lower() for value in inputs.get("files") or [])
    args = _arguments(inputs)
    arg_text = " ".join(str(value).lower() for value in args.values() if isinstance(value, str))
    haystack = " ".join([request_text, file_text, arg_text]).strip()
    if not haystack:
        return None
    best: tuple[int, dict[str, Any] | None] = (0, None)
    for planlet in planlets:
        score = _planlet_score(planlet, haystack)
        if score > best[0]:
            best = (score, planlet)
    return best[1] if best[0] >= 2 else None


def _planlet_score(planlet: dict[str, Any], haystack: str) -> int:
    score = 0
    applicability = planlet.get("applicability") if isinstance(planlet.get("applicability"), dict) else {}
    for ext in applicability.get("file_extensions") or []:
        if str(ext).lower() and str(ext).lower() in haystack:
            score += 3
    for keyword in planlet.get("keywords") or []:
        word = str(keyword).lower()
        if len(word) >= 3 and word in haystack:
            score += 1
    planlet_id = str(planlet.get("id") or "").lower().replace("_", " ")
    for word in re.findall(r"[a-z0-9]{3,}", planlet_id):
        if word in haystack:
            score += 1
    return score


def _missing_planlet_bindings(planlet: dict[str, Any], inputs: dict[str, Any]) -> list[str]:
    bindings = planlet.get("bindings") if isinstance(planlet.get("bindings"), dict) else {}
    required = [str(name) for name in bindings.get("required") or [] if name]
    args = _arguments(inputs)
    missing = []
    for name in required:
        if not _has_bound_value(_lookup_planlet_binding(name, args, inputs, planlet)):
            missing.append(name)
    return missing


def _execute_planlet(
    spec: dict[str, Any],
    planlet: dict[str, Any],
    inputs: dict[str, Any],
) -> dict[str, Any]:
    step_results: dict[str, Any] = {}
    debug_steps: list[dict[str, Any]] = []
    operators = [op for op in spec.get("operators") or [] if isinstance(op, dict)]
    args = _arguments(inputs)
    for step in [item for item in planlet.get("steps") or [] if isinstance(item, dict)]:
        operation = str(step.get("operation") or "").strip()
        if not operation:
            continue
        operator = _operator_by_name(operators, operation)
        if operator is None:
            if bool(step.get("optional")):
                continue
            return _blocked(
                spec,
                "planlet_operator_missing",
                f"Planlet step references an unknown operation: {operation}",
                {"planlet": planlet.get("id"), "step": step.get("id")},
            )
        step_args = _bind_planlet_step_arguments(step, args, inputs, planlet)
        step_inputs = dict(inputs)
        step_inputs["operation"] = operation
        step_inputs["arguments"] = step_args
        step_inputs["execute"] = True
        validation = _validate_operator_arguments(operator, step_inputs)
        if validation["errors"]:
            if bool(step.get("optional")):
                debug_steps.append(
                    {
                        "step": step.get("id") or operation,
                        "operation": operation,
                        "status": "skipped",
                        "validation": validation,
                    }
                )
                continue
            return _blocked(
                spec,
                "invalid_planlet_arguments",
                f"Invalid arguments for planlet step {step.get('id') or operation}.",
                {
                    "planlet": planlet.get("id"),
                    "step": step.get("id"),
                    "operator": operation,
                    **validation,
                },
            )
        execution = _execute_operator(spec, operator, step_inputs)
        key = str(step.get("id") or operation)
        step_results[key] = {
            "operator": operation,
            "executor": execution.get("executor"),
            "returncode": execution.get("returncode"),
            "stdout": execution.get("stdout"),
            "stderr": execution.get("stderr"),
            "output": execution.get("final_output"),
        }
        debug_steps.append(
            {
                "step": key,
                "operation": operation,
                "status": execution.get("status"),
            }
        )
        output = execution.get("final_output")
        if isinstance(output, dict):
            for output_key in (
                "output_path",
                "output_pptx",
                "output_docx",
                "output_file",
                "path",
            ):
                if _has_bound_value(output.get(output_key)):
                    args["last_output_path"] = output[output_key]
                    break
        if execution.get("status") != "completed":
            if bool(step.get("optional")):
                continue
            return _blocked(
                spec,
                execution.get("deopt_reason") or "planlet_step_failed",
                execution.get("error") or "Planlet step failed.",
                {
                    "planlet": planlet.get("id"),
                    "step": key,
                    "operator": operation,
                    "execution": execution,
                },
            )
    boundary = str(spec.get("boundary") or "adapter")
    return _completed(
        spec,
        runtime_status="planlet_executed",
        contribution_type="execution" if boundary in _SOLVER_BOUNDARIES else "typed_evidence",
        final_output={
            "planlet": _planlet_summary(planlet),
            "step_results": step_results,
            "stop_policy": planlet.get("stop_policy") or {},
        },
        selected_operator=None,
        selected_planlet=planlet,
        requires_agent_continuation=boundary not in _SOLVER_BOUNDARIES,
        step_results=step_results,
        debug_trace={"planlet_steps": debug_steps},
    )


def _bind_planlet_step_arguments(
    step: dict[str, Any],
    args: dict[str, Any],
    inputs: dict[str, Any],
    planlet: dict[str, Any],
) -> dict[str, Any]:
    step_spec = step.get("arguments") if isinstance(step.get("arguments"), dict) else {}
    if step_spec.get("*") == "${arguments}":
        return dict(args)
    bound: dict[str, Any] = {}
    for name, value in step_spec.items():
        resolved = _resolve_planlet_value(value, args, inputs, planlet)
        if resolved is _MISSING:
            continue
        bound[name] = resolved
    return bound


_MISSING = object()


def _resolve_planlet_value(
    value: Any,
    args: dict[str, Any],
    inputs: dict[str, Any],
    planlet: dict[str, Any],
) -> Any:
    if not isinstance(value, str) or not value.startswith("${") or not value.endswith("}"):
        return value
    expr = value[2:-1]
    name, default = _split_planlet_default(expr)
    if name == "arguments":
        return dict(args)
    if name.startswith("arguments."):
        return args.get(name.removeprefix("arguments."), default)
    resolved = _lookup_planlet_binding(name, args, inputs, planlet)
    if _has_bound_value(resolved):
        return resolved
    if default is not _MISSING:
        return default
    return _MISSING


def _split_planlet_default(expr: str) -> tuple[str, Any]:
    if ":" not in expr:
        return expr, _MISSING
    name, raw_default = expr.split(":", 1)
    raw_default = raw_default.strip()
    if raw_default == "":
        return name, ""
    if raw_default.lower() == "true":
        return name, True
    if raw_default.lower() == "false":
        return name, False
    if raw_default == "[]":
        return name, []
    if raw_default == "{}":
        return name, {}
    return name, raw_default


def _lookup_planlet_binding(
    name: str,
    args: dict[str, Any],
    inputs: dict[str, Any],
    planlet: dict[str, Any],
) -> Any:
    if _has_bound_value(args.get(name)):
        return args.get(name)
    bindings = planlet.get("bindings") if isinstance(planlet.get("bindings"), dict) else {}
    autobind = bindings.get("autobind") if isinstance(bindings.get("autobind"), dict) else {}
    for alias in autobind.get(name) or []:
        value = _lookup_binding_alias(str(alias), args, inputs)
        if _has_bound_value(value):
            return value
    for alias in _argument_aliases(name, ""):
        value = args.get(alias)
        if _has_bound_value(value):
            return value
    return None


def _lookup_binding_alias(alias: str, args: dict[str, Any], inputs: dict[str, Any]) -> Any:
    if alias == "files[0]":
        files = inputs.get("files")
        if isinstance(files, list) and files:
            return files[0]
        return None
    return args.get(alias)


def _operator_by_name(operators: list[dict[str, Any]], name: str) -> dict[str, Any] | None:
    requested = name.lower()
    for op in operators:
        if requested in {
            str(op.get("id") or "").lower(),
            str(op.get("name") or "").lower(),
        }:
            return op
    return None


def _planlet_guidance(planlet: dict[str, Any], missing: list[str]) -> dict[str, Any]:
    return {
        "selected_planlet": _planlet_summary(planlet),
        "missing_bindings": missing,
        "required_bindings": list((planlet.get("bindings") or {}).get("required") or []),
        "optional_bindings": list((planlet.get("bindings") or {}).get("optional") or []),
        "steps": [
            {
                "id": str(step.get("id") or ""),
                "operation": str(step.get("operation") or ""),
                "optional": bool(step.get("optional", False)),
            }
            for step in planlet.get("steps") or []
            if isinstance(step, dict)
        ],
        "hint": "Provide the missing bindings under `arguments`, or call an exact operation directly.",
    }


def _validate_operator_arguments(
    operator: dict[str, Any],
    inputs: dict[str, Any],
) -> dict[str, Any]:
    args = _arguments(inputs, operator)
    errors: list[dict[str, Any]] = []
    for field in operator.get("inputs") or []:
        if not isinstance(field, dict):
            continue
        name = str(field.get("name") or "")
        if not name:
            continue
        value = args.get(name)
        alias_bound = any(
            _has_bound_value(args.get(alias))
            for alias in _argument_aliases(name, str(operator.get("kind") or ""))
        )
        if field.get("required", False) and not _has_bound_value(value) and not alias_bound:
            errors.append(
                {
                    "field": name,
                    "reason": "missing_required",
                    "accepted_aliases": _argument_aliases(
                        name, str(operator.get("kind") or "")
                    ),
                }
            )
            continue
        if _has_bound_value(value):
            type_error = _argument_type_error(name, value, str(field.get("type") or ""))
            if type_error:
                errors.append(type_error)
    return {
        "errors": errors,
        "bound_arguments": {
            key: value
            for key, value in args.items()
            if isinstance(value, (str, int, float, bool, list, dict)) or value is None
        },
        "operation": operator.get("name") or operator.get("id"),
    }


def _select_operator(spec: dict[str, Any], inputs: dict[str, Any]) -> dict[str, Any] | None:
    operators = [op for op in spec.get("operators") or [] if isinstance(op, dict)]
    if not operators:
        return None
    requested = str(inputs.get("operation") or inputs.get("topic") or "").strip().lower()
    if requested:
        for op in operators:
            names = {
                str(op.get("id") or "").lower(),
                str(op.get("name") or "").lower(),
            }
            if requested in names:
                return op
        for op in operators:
            if requested in str(op.get("name") or "").lower():
                return op
    request = " ".join(
        str(value)
        for value in [inputs.get("request"), inputs.get("task"), inputs.get("topic")]
        if value
    ).lower()
    if request:
        best: tuple[int, dict[str, Any] | None] = (0, None)
        for op in operators:
            words = set(str(item).lower() for item in op.get("keywords") or [])
            words.update(re.findall(r"[a-z0-9_]{3,}", str(op.get("name") or "").lower()))
            score = sum(1 for word in words if word and word in request)
            if score > best[0]:
                best = (score, op)
        if best[1] is not None:
            return best[1]
    return operators[0] if len(operators) == 1 else None


def _network_blocked(
    spec: dict[str, Any],
    operator: dict[str, Any],
    inputs: dict[str, Any],
) -> bool:
    if _is_lossless_reference_operator(operator):
        return False
    if _truthy(inputs.get("allow_network")):
        return False
    policy = spec.get("execution_policy") or {}
    if bool(policy.get("allow_network_by_default")):
        return False
    flags = set(operator.get("risk_flags") or [])
    if not flags:
        flags = set(spec.get("risk_flags") or [])
    return "network" in flags


def _completed(
    spec: dict[str, Any],
    *,
    runtime_status: str,
    contribution_type: str,
    final_output: Any,
    requires_agent_continuation: bool,
    selected_operator: dict[str, Any] | None = None,
    selected_planlet: dict[str, Any] | None = None,
    step_results: dict[str, Any] | None = None,
    debug_trace: dict[str, Any] | None = None,
) -> dict[str, Any]:
    boundary = str(spec.get("boundary") or "adapter")
    return {
        "status": "completed",
        "runtime_status": runtime_status,
        "boundary": boundary,
        "can_produce_final_answer": boundary in _SOLVER_BOUNDARIES,
        "requires_agent_continuation": requires_agent_continuation,
        "contribution_type": contribution_type,
        "final_output": final_output,
        "produced_outputs": [],
        "selected_operator": _operator_summary(selected_operator),
        "selected_planlet": _planlet_summary(selected_planlet),
        "step_results": step_results or {},
        "artifact_boundary": _summary(spec),
        "debug_trace": debug_trace or {},
    }


def _blocked(
    spec: dict[str, Any],
    reason: str,
    message: str,
    details: dict[str, Any] | None = None,
) -> dict[str, Any]:
    boundary = str(spec.get("boundary") or "adapter")
    return {
        "status": "blocked",
        "runtime_status": "blocked",
        "boundary": boundary,
        "can_produce_final_answer": boundary in _SOLVER_BOUNDARIES,
        "requires_agent_continuation": True,
        "contribution_type": "none",
        "final_output": None,
        "produced_outputs": [],
        "deopt_reason": reason,
        "errors": [{"reason": reason, "message": message, "details": details or {}}],
        "block": {
            "message": message,
            "details": details or {},
            "suggested_actions": [
                "Continue with normal agent tools or use a more specific operator."
            ],
        },
        "artifact_boundary": _summary(spec),
    }


def _inputs(state: dict[str, Any]) -> dict[str, Any]:
    global_inputs = dict(state.get("global_inputs") or {})
    nested = global_inputs.pop("inputs", None)
    if isinstance(nested, dict):
        merged = dict(nested)
        merged.update(global_inputs)
        return merged
    return global_inputs


def _arguments(inputs: dict[str, Any], operator: dict[str, Any] | None = None) -> dict[str, Any]:
    args = inputs.get("arguments")
    merged = dict(args) if isinstance(args, dict) else {}
    _bind_context_file_if_unambiguous(merged, inputs, operator or {})
    _coerce_operator_argument_aliases(merged, operator or {})
    return merged


def _bind_context_file_if_unambiguous(
    args: dict[str, Any],
    inputs: dict[str, Any],
    operator: dict[str, Any],
) -> None:
    files = inputs.get("files")
    if not isinstance(files, list) or not files:
        return
    first_file = files[0]
    if not isinstance(first_file, str) or not first_file:
        return
    missing_path_args = [
        name
        for name in _required_operator_inputs(operator)
        if not _has_bound_value(args.get(name)) and _is_input_path_like(name)
    ]
    if len(missing_path_args) == 1:
        args[missing_path_args[0]] = first_file


def _coerce_operator_argument_aliases(args: dict[str, Any], operator: dict[str, Any]) -> None:
    if not operator:
        return
    kind = str(operator.get("kind") or "")
    for field in operator.get("inputs") or []:
        if not isinstance(field, dict):
            continue
        name = str(field.get("name") or "")
        if not name or _has_bound_value(args.get(name)):
            continue
        if kind == "pptx_update_text_boxes" and _normalize_argument_name(name) == "edits":
            continue
        for alias in _argument_aliases(name, kind):
            if alias == name:
                continue
            value = args.get(alias)
            if _has_bound_value(value):
                args[name] = value
                break


def _argument_aliases(name: str, kind: str) -> list[str]:
    normalized = _normalize_argument_name(name)
    aliases: dict[str, list[str]] = {
        "path": ["path", "file", "input", "input_path", "docx_path", "pptx_path", "bibtex"],
        "file": ["file", "path", "input", "input_path", "bibtex", "bibtex_file"],
        "input": ["input", "input_path", "path", "file"],
        "input_file": ["input_file", "file", "path", "input", "input_path"],
        "input_path": [
            "input_path",
            "input",
            "path",
            "file",
            "source",
            "source_path",
            "docx_path",
            "pptx_path",
            "input_pptx",
        ],
        "input_pptx": ["input_pptx", "pptx_path", "input_path", "path", "input", "file"],
        "output": ["output", "output_path", "output_file", "destination", "dest", "answer"],
        "output_file": ["output_file", "output", "output_path", "destination", "dest", "answer"],
        "output_path": [
            "output_path",
            "output",
            "output_file",
            "destination",
            "dest",
            "answer",
            "answer_file",
            "output_docx",
            "output_pptx",
        ],
        "output_pptx": [
            "output_pptx",
            "output_path",
            "output",
            "output_file",
            "destination",
            "dest",
        ],
        "replacements": ["replacements", "replacement_map", "values", "mapping"],
        "markers": ["markers", "marker_specs", "blocks"],
        "items": ["items", "references", "titles"],
        "edits": ["edits", "updates", "targets", "text_boxes"],
    }
    result = list(aliases.get(normalized, [normalized]))
    if kind == "citation_validate_local" and normalized == "file":
        result.extend(["bib", "bib_file", "test_bib"])
    return _dedupe_strings(result)


def _dedupe_strings(items: list[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for item in items:
        if item in seen:
            continue
        seen.add(item)
        result.append(item)
    return result


def _argument_type_error(name: str, value: Any, expected_type: str) -> dict[str, Any] | None:
    expected = expected_type.strip().lower()
    if expected in {"", "any"}:
        return None
    if expected in {"object", "dict", "mapping"}:
        if isinstance(value, dict) or _string_parses_as(value, dict):
            return None
        return {"field": name, "reason": "type_mismatch", "expected": "object"}
    if expected in {"array", "list"}:
        if isinstance(value, list) or _string_parses_as(value, list):
            return None
        return {"field": name, "reason": "type_mismatch", "expected": "array"}
    if expected in {"boolean", "bool"}:
        if isinstance(value, bool) or str(value).strip().lower() in {
            "true",
            "false",
            "1",
            "0",
            "yes",
            "no",
            "on",
            "off",
        }:
            return None
        return {"field": name, "reason": "type_mismatch", "expected": "boolean"}
    if expected in {"integer", "int"}:
        try:
            int(value)
            return None
        except Exception:
            return {"field": name, "reason": "type_mismatch", "expected": "integer"}
    if expected in {"number", "float"}:
        try:
            float(value)
            return None
        except Exception:
            return {"field": name, "reason": "type_mismatch", "expected": "number"}
    if expected in {"string", "str"} and not isinstance(value, str):
        return {"field": name, "reason": "type_mismatch", "expected": "string"}
    return None


def _string_parses_as(value: Any, expected_type: type) -> bool:
    if not isinstance(value, str) or not value.strip():
        return False
    try:
        parsed = json.loads(value)
    except Exception:
        try:
            parsed = ast.literal_eval(value)
        except Exception:
            return False
    return isinstance(parsed, expected_type)


def _has_bound_value(value: Any) -> bool:
    if value is None:
        return False
    if isinstance(value, str) and not value.strip():
        return False
    return True


def _required_operator_inputs(operator: dict[str, Any]) -> list[str]:
    names: list[str] = []
    for field in operator.get("inputs") or []:
        if not isinstance(field, dict) or not field.get("required", False):
            continue
        name = str(field.get("name") or "")
        if name:
            names.append(name)
    return names


def _is_input_path_like(name: str) -> bool:
    normalized = _normalize_argument_name(name)
    if _is_output_path_like(normalized):
        return False
    return (
        normalized
        in {
            "file",
            "input",
            "input_file",
            "input_path",
            "source",
            "source_file",
            "source_path",
            "path",
            "config",
            "config_file",
            "config_path",
            "query_file",
            "data_file",
            "csv",
            "csv_file",
            "json_file",
            "pdf",
            "pdf_file",
            "image",
            "image_file",
            "video",
            "video_file",
            "audio",
            "audio_file",
            "bibtex",
            "bibtex_file",
        }
        or normalized.startswith("input_")
        or normalized.endswith("_file")
        or normalized.endswith("_path")
    )


def _is_output_path_like(name: str) -> bool:
    normalized = _normalize_argument_name(name)
    return (
        normalized
        in {
            "o",
            "out",
            "output",
            "output_file",
            "output_path",
            "output_dir",
            "output_directory",
            "output_prefix",
            "report",
            "report_file",
            "answer",
            "answer_file",
            "destination",
            "destination_file",
            "dest",
            "dest_file",
        }
        or normalized.startswith("output_")
        or normalized.endswith("_output")
    )


def _normalize_argument_name(name: str) -> str:
    return re.sub(r"[^a-z0-9_]+", "_", name.strip().lower().replace("-", "_")).strip("_")


def _safe_int(value: Any, *, default: int) -> int:
    try:
        return int(value)
    except Exception:
        return default


def _sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        while True:
            chunk = f.read(1024 * 1024)
            if not chunk:
                break
            h.update(chunk)
    return h.hexdigest()


def _required_workspace_path(project_dir: Path, args: dict[str, Any], name: str) -> Path:
    value = args.get(name)
    if not isinstance(value, str) or not value.strip():
        raise ValueError(f"Missing required argument: {name}")
    return _resolve_workspace_path(project_dir, value.strip())


def _required_workspace_path_alias(
    project_dir: Path,
    args: dict[str, Any],
    primary: str,
    *aliases: str,
) -> Path:
    for name in (primary, *aliases):
        value = args.get(name)
        if isinstance(value, str) and value.strip():
            return _resolve_workspace_path(project_dir, value.strip())
    expected = ", ".join((primary, *aliases))
    raise ValueError(f"Missing required argument: {primary} (accepted aliases: {expected})")


def _optional_workspace_path(project_dir: Path, args: dict[str, Any], name: str) -> Path | None:
    value = args.get(name)
    if not isinstance(value, str) or not value.strip():
        return None
    return _resolve_workspace_path(project_dir, value.strip())


def _mapping_arg(
    args: dict[str, Any],
    name: str,
    *,
    required: bool = True,
) -> dict[str, Any]:
    value = args.get(name)
    if isinstance(value, str) and value.strip():
        try:
            value = json.loads(value)
        except Exception:
            value = ast.literal_eval(value)
    if value is None and not required:
        return {}
    if not isinstance(value, dict):
        raise ValueError(f"Missing required object argument: {name}")
    return value


def _list_arg(
    args: dict[str, Any],
    name: str,
    *,
    required: bool = False,
) -> list[Any]:
    value = args.get(name)
    if isinstance(value, str) and value.strip():
        try:
            value = json.loads(value)
        except Exception:
            value = ast.literal_eval(value)
    if value is None:
        if required:
            raise ValueError(f"Missing required list argument: {name}")
        return []
    if not isinstance(value, list):
        raise ValueError(f"Argument must be a list: {name}")
    return value


def _truthy_arg(args: dict[str, Any], name: str, *, default: bool) -> bool:
    if name not in args:
        return default
    return _truthy(args.get(name))


def _docx_paragraphs(
    doc: Any,
    *,
    include_tables: bool,
    include_headers_footers: bool,
) -> list[Any]:
    paragraphs = list(doc.paragraphs)
    if include_tables:
        _extend_docx_table_paragraphs(paragraphs, doc.tables)
    if include_headers_footers:
        for section in doc.sections:
            paragraphs.extend(section.header.paragraphs)
            paragraphs.extend(section.footer.paragraphs)
            if include_tables:
                _extend_docx_table_paragraphs(paragraphs, section.header.tables)
                _extend_docx_table_paragraphs(paragraphs, section.footer.tables)
    return paragraphs


def _extend_docx_table_paragraphs(paragraphs: list[Any], tables: Any) -> None:
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                paragraphs.extend(cell.paragraphs)
                _extend_docx_table_paragraphs(paragraphs, cell.tables)


def _docx_table_text_items(tables: Any, *, prefix: str = "") -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for table_idx, table in enumerate(tables):
        table_path = f"{prefix}{table_idx}" if not prefix else f"{prefix}.{table_idx}"
        for row_idx, row in enumerate(table.rows):
            for cell_idx, cell in enumerate(row.cells):
                text = "\n".join(p.text for p in cell.paragraphs if p.text)
                if text:
                    items.append(
                        {
                            "table": table_path,
                            "row": row_idx,
                            "cell": cell_idx,
                            "text": text,
                        }
                    )
                items.extend(_docx_table_text_items(cell.tables, prefix=table_path))
    return items


def _replace_paragraph_text(paragraph: Any, replacements: dict[str, Any]) -> int:
    original = "".join(run.text for run in paragraph.runs) if paragraph.runs else paragraph.text
    updated = original
    for needle, replacement in replacements.items():
        updated = updated.replace(str(needle), str(replacement))
    if updated == original:
        return 0
    if paragraph.runs:
        paragraph.runs[0].text = updated
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = updated
    return 1


def _marker_specs(args: dict[str, Any]) -> list[dict[str, Any]]:
    markers = args.get("markers")
    if isinstance(markers, str) and markers.strip():
        try:
            markers = json.loads(markers)
        except Exception:
            markers = ast.literal_eval(markers)
    if isinstance(markers, dict):
        markers = [markers]
    if not markers:
        start = args.get("start") or args.get("start_marker")
        end = args.get("end") or args.get("end_marker")
        if start and end:
            return [
                {
                    "name": args.get("name") or "marker",
                    "start": start,
                    "end": end,
                    "keep": args.get("keep"),
                }
            ]
    if not isinstance(markers, list) or not all(isinstance(item, dict) for item in markers):
        raise ValueError("Missing required argument: markers")
    return markers


def _apply_pptx_positioning(shape: Any, prs: Any, edit: dict[str, Any]) -> None:
    position = str(edit.get("position") or "").strip().lower().replace("-", "_")
    if not position:
        return
    from pptx.util import Inches

    if position == "bottom_center":
        margin = int(edit.get("margin") or Inches(0.25))
        bottom_margin = int(edit.get("bottom_margin") or Inches(0.18))
        height = int(edit.get("height") or min(int(shape.height), int(Inches(0.45))))
        shape.left = margin
        shape.width = int(prs.slide_width) - (2 * margin)
        shape.height = height
        shape.top = int(prs.slide_height) - bottom_margin - height
        return
    if position == "center":
        shape.left = (int(prs.slide_width) - int(shape.width)) // 2
        shape.top = (int(prs.slide_height) - int(shape.height)) // 2
        return
    if position == "bottom":
        bottom_margin = int(edit.get("bottom_margin") or Inches(0.18))
        shape.top = int(prs.slide_height) - bottom_margin - int(shape.height)


def _apply_pptx_text_frame_options(shape: Any, edit: dict[str, Any]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    from pptx.enum.text import MSO_AUTO_SIZE

    text_frame = shape.text_frame
    if "word_wrap" in edit:
        text_frame.word_wrap = _bool_value(edit["word_wrap"])
    if _truthy_arg(edit, "fit_to_one_line", default=False):
        text_frame.word_wrap = False
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        if attr in edit:
            setattr(text_frame, attr, int(edit[attr]))
    if _truthy_arg(edit, "zero_margins", default=False) or _truthy_arg(
        edit, "fit_to_one_line", default=False
    ):
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0


def _apply_pptx_font_style(
    shape: Any, edit: dict[str, Any], rgb_color: Any, pt: Any, pp_align: Any
) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    color = _pptx_rgb(edit.get("font_color"), rgb_color) if "font_color" in edit else None
    alignment = _pptx_alignment(edit.get("alignment", edit.get("align")), pp_align)
    for paragraph in shape.text_frame.paragraphs:
        if alignment is not None:
            paragraph.alignment = alignment
        if not paragraph.runs and paragraph.text:
            text = paragraph.text
            paragraph.clear()
            paragraph.add_run().text = text
        for run in paragraph.runs:
            if "font_name" in edit:
                run.font.name = str(edit["font_name"])
            if "font_size" in edit:
                run.font.size = pt(float(edit["font_size"]))
            if "bold" in edit:
                run.font.bold = _bool_value(edit["bold"])
            if color is not None:
                run.font.color.rgb = color


def _pptx_rgb(value: Any, rgb_color: Any) -> Any:
    if isinstance(value, (list, tuple)) and len(value) == 3:
        return rgb_color(int(value[0]), int(value[1]), int(value[2]))
    raw = str(value).strip()
    if raw.startswith("#"):
        raw = raw[1:]
    if len(raw) != 6:
        raise ValueError(f"Invalid PPTX RGB color: {value!r}")
    return rgb_color(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _pptx_alignment(value: Any, pp_align: Any) -> Any | None:
    if value is None:
        return None
    mapping = {
        "center": pp_align.CENTER,
        "left": pp_align.LEFT,
        "right": pp_align.RIGHT,
        "justify": pp_align.JUSTIFY,
    }
    return mapping.get(str(value).strip().lower())


def _bool_value(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    return str(value).strip().lower() in {"1", "true", "yes", "y", "on"}


def _set_pptx_auto_number(paragraph: Any) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    paragraph.level = 0
    ppr = paragraph._p.get_or_add_pPr()
    for child in list(ppr):
        if child.tag in {qn("a:buNone"), qn("a:buChar"), qn("a:buAutoNum")}:
            ppr.remove(child)
    bu_auto_num = OxmlElement("a:buAutoNum")
    bu_auto_num.set("type", "arabicPeriod")
    ppr.insert(0, bu_auto_num)


def _pptx_shape_matches(shape: Any, edit: dict[str, Any]) -> bool:
    has_selector = False
    if "shape_id" in edit:
        has_selector = True
        if int(edit["shape_id"]) != int(shape.shape_id):
            return False
    if "shape_name" in edit:
        has_selector = True
        if str(edit["shape_name"]) != str(shape.name):
            return False
    text = shape.text if getattr(shape, "has_text_frame", False) else ""
    if "match_text" in edit:
        has_selector = True
        if str(edit["match_text"]) != text:
            return False
    if "contains" in edit:
        has_selector = True
        if str(edit["contains"]) not in text:
            return False
    return has_selector


def _timeout(spec: dict[str, Any]) -> int:
    raw = (spec.get("execution_policy") or {}).get("timeout_sec", 20)
    try:
        return max(1, int(raw))
    except Exception:
        return 20


def _python_executable(project_dir: Path) -> str:
    host_python = project_dir / "host_python"
    if host_python.exists():
        return str(host_python)
    return sys.executable


def _resolve_workspace_path(project_dir: Path, value: str) -> Path:
    path = Path(value)
    if path.is_absolute():
        raw = path.as_posix()
        mappings = {
            "/root": project_dir,
            "/app": project_dir,
            "/data": project_dir / "data",
            "/output": project_dir / "output",
        }
        for prefix, target in mappings.items():
            if raw == prefix:
                return target
            if raw.startswith(prefix + "/"):
                return target / raw.removeprefix(prefix + "/")
        return path
    return project_dir / path


def _bind_command(command: str, args: dict[str, Any]) -> str:
    bound = command
    for key, value in args.items():
        bound = bound.replace(f"<{key}>", shlex.quote(str(value)))
        bound = bound.replace(f"{{{key}}}", shlex.quote(str(value)))
    return bound


def _parse_output(stdout: str) -> Any:
    text = stdout.strip()
    if not text:
        return ""
    try:
        return json.loads(text)
    except Exception:
        pass
    try:
        return ast.literal_eval(text)
    except Exception:
        return text


def _operator_guidance(operator: dict[str, Any]) -> str:
    parts = [
        str(operator.get("name") or operator.get("id") or ""),
        str(operator.get("description") or operator.get("summary") or ""),
    ]
    if operator.get("script_path"):
        parts.append(f"script: {operator['script_path']}")
    if operator.get("commands"):
        parts.append("commands:\n" + "\n".join(str(item) for item in operator["commands"]))
    return "\n\n".join(part for part in parts if part).strip()


def _operator_names(spec: dict[str, Any]) -> list[str]:
    return [
        str(op.get("name") or op.get("id"))
        for op in spec.get("operators") or []
        if isinstance(op, dict)
    ]


def _operator_summary(operator: dict[str, Any] | None) -> dict[str, Any]:
    if not operator:
        return {}
    return {
        "id": operator.get("id"),
        "name": operator.get("name"),
        "kind": operator.get("kind"),
        "risk_flags": operator.get("risk_flags", []),
    }


def _planlet_summary(planlet: dict[str, Any] | None) -> dict[str, Any]:
    if not planlet:
        return {}
    bindings = planlet.get("bindings") if isinstance(planlet.get("bindings"), dict) else {}
    return {
        "id": planlet.get("id"),
        "intent": planlet.get("intent"),
        "required_bindings": list(bindings.get("required") or []),
        "step_count": len(planlet.get("steps") or []),
    }


def _summary(spec: dict[str, Any]) -> dict[str, Any]:
    boundary = str(spec.get("boundary") or "adapter")
    selection = spec.get("selection_policy") if isinstance(spec.get("selection_policy"), dict) else {}
    planlets = [item for item in spec.get("planlets") or [] if isinstance(item, dict)]
    return {
        "artifact_id": spec.get("artifact_id", ""),
        "skill_name": spec.get("skill_name", ""),
        "boundary": boundary,
        "direct_call_level": selection.get("direct_call_level", ""),
        "recommended_planlets": [str(item.get("id") or "") for item in planlets[:3]],
        "planlet_count": len(planlets),
        "can_produce_final_answer": boundary in _SOLVER_BOUNDARIES,
        "requires_agent_continuation": boundary not in _SOLVER_BOUNDARIES,
        "risk_flags": list(spec.get("risk_flags") or []),
    }


def _truthy(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "on"}
    return bool(value)
