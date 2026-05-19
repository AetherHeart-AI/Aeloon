from __future__ import annotations

import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

import pytest

from aeloon.plugins.SkillGraph.skillgraph.boundary import build_boundary_spec
from aeloon.plugins.SkillGraph.skillgraph.boundary_runtime import run_boundary_artifact
from aeloon.plugins.SkillGraph.skillgraph.dispatcher_codegen import extract_dispatcher_capabilities
from aeloon.plugins.SkillGraph.skillgraph.models import RuntimeManifest
from aeloon.plugins.SkillGraph.skillgraph.package import build_skill_package


def _write_lossless_sandbox(tmp_path: Path) -> tuple[Path, dict]:
    sandbox = tmp_path / "artifact.sandbox"
    skill_root = sandbox / "skill"
    skill_root.mkdir(parents=True)
    (skill_root / "SKILL.md").write_text("# Skill\n\nUse this skill for tests.\n", encoding="utf-8")
    assets = [
        {
            "path": "SKILL.md",
            "asset_type": "instruction",
            "is_text": True,
            "size_bytes": (skill_root / "SKILL.md").stat().st_size,
            "sha256": "",
            "line_count": 3,
        }
    ]
    spec = {
        "boundary": "adapter",
        "risk_flags": ["network"],
        "scope": {
            "lossless_reference": {
                "strict_validated": False,
                "package_hash": "pkg",
                "capsule_hash": "cap",
                "asset_count": 1,
                "text_asset_count": 1,
                "binary_asset_count": 0,
                "assets": assets,
            }
        },
        "operators": [
            {"id": "list_skill_assets", "name": "list_skill_assets", "kind": "skill_asset_list"},
            {"id": "get_skill_asset", "name": "get_skill_asset", "kind": "skill_asset_get"},
            {"id": "search_skill_docs", "name": "search_skill_docs", "kind": "skill_doc_search"},
        ],
        "execution_policy": {"allow_network_by_default": False},
    }
    return sandbox, spec


def test_lossless_reference_ops_do_not_inherit_artifact_network_risk(tmp_path: Path) -> None:
    sandbox, spec = _write_lossless_sandbox(tmp_path)

    cases = [
        ("list_skill_assets", {}),
        ("get_skill_asset", {"path": "SKILL.md", "max_chars": 20}),
        ("search_skill_docs", {"query": "skill", "max_results": 1}),
    ]
    for operation, arguments in cases:
        result = run_boundary_artifact(
            spec,
            {
                "global_inputs": {
                    "sandbox_dir": str(sandbox),
                    "operation": operation,
                    "arguments": arguments,
                },
                "step_results": {},
            },
        )
        assert result["status"] == "completed", operation
        assert result["runtime_status"] == "guidance_only"


def test_get_skill_asset_requires_explicit_skill_asset_path(tmp_path: Path) -> None:
    sandbox, spec = _write_lossless_sandbox(tmp_path)

    result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "sandbox_dir": str(sandbox),
                "operation": "get_skill_asset",
                "arguments": {},
                "files": [str(tmp_path / "workspace_template.docx")],
            },
            "step_results": {},
        },
    )

    assert result["status"] == "blocked"
    assert "Missing required argument: path" in result["block"]["message"]


def test_script_asset_without_schema_uses_context_file_and_parses_python_literal(
    tmp_path: Path,
) -> None:
    sandbox = tmp_path / "artifact.sandbox"
    script_dir = sandbox / "skill" / "scripts"
    script_dir.mkdir(parents=True)
    script = script_dir / "mesh_tool.py"
    script.write_text(
        "import sys\nprint({'argc': len(sys.argv), 'arg': sys.argv[1] if len(sys.argv) > 1 else ''})\n",
        encoding="utf-8",
    )
    input_file = tmp_path / "scan_data.stl"
    input_file.write_text("mesh", encoding="utf-8")
    spec = {
        "boundary": "typed_operator",
        "operators": [
            {
                "id": "mesh_tool",
                "name": "mesh_tool",
                "kind": "script_asset",
                "script_path": "scripts/mesh_tool.py",
                "inputs": [],
                "risk_flags": [],
            }
        ],
        "execution_policy": {"timeout_sec": 10},
    }

    result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "sandbox_dir": str(sandbox),
                "operation": "mesh_tool",
                "files": [str(input_file)],
            },
            "step_results": {},
        },
    )

    assert result["status"] == "completed"
    assert result["final_output"]["argc"] == 2
    assert result["final_output"]["arg"] == str(input_file)


def test_docx_skill_only_capabilities_promote_to_typed_operator(tmp_path: Path) -> None:
    skill_root = tmp_path / "docx"
    skill_root.mkdir()
    entry = skill_root / "SKILL.md"
    entry.write_text(
        "# DOCX\n\nUse python-docx for split placeholders.\n\n```bash\npython ad_hoc.py\n```\n",
        encoding="utf-8",
    )
    package = build_skill_package(skill_root)

    capabilities = extract_dispatcher_capabilities(package, entry)
    names = {cap.name for cap in capabilities}
    kinds = {cap.kind for cap in capabilities}

    assert "replace_docx_text" in names
    assert "validate_docx_text" in names
    assert "shell_command" not in kinds

    spec = build_boundary_spec(
        strategy="dispatcher",
        skill_name=package.slug,
        description="DOCX",
        runtime_manifest=RuntimeManifest(),
        package=package,
        metadata={
            "task_context": {},
            "compilability": {
                "strategy": "dispatcher",
                "kind": "toolkit_dispatcher",
                "source_shape": "toolkit",
                "confidence": "medium",
            },
        },
        dispatcher_capabilities=capabilities,
    )

    assert spec.boundary == "typed_operator"
    assert spec.selection_policy["direct_call_level"] == "operator_first"
    planlet_ids = {planlet.id for planlet in spec.planlets}
    assert "docx_template_fill_and_validate" in planlet_ids
    assert spec.scope["kind"] == "skill_package_generic"
    assert not spec.scope["dynamic_sources_used_for_codegen"]


def test_docx_replace_and_validate_generic_operator(tmp_path: Path) -> None:
    pytest.importorskip("docx")
    from docx import Document

    source = tmp_path / "template.docx"
    output = tmp_path / "filled.docx"
    doc = Document()
    paragraph = doc.add_paragraph()
    paragraph.add_run("Hello {{")
    paragraph.add_run("NAME")
    paragraph.add_run("}}")
    table = doc.add_table(rows=1, cols=1)
    nested = table.cell(0, 0).add_table(rows=1, cols=1)
    nested.cell(0, 0).text = "Role: {{ROLE}}"
    doc.save(source)
    spec = {
        "boundary": "typed_operator",
        "operators": [
            {
                "id": "replace_docx_text",
                "name": "replace_docx_text",
                "kind": "docx_replace_text",
                "inputs": [
                    {"name": "input_path", "required": True},
                    {"name": "output_path", "required": True},
                    {"name": "replacements", "required": True, "type": "object"},
                ],
            },
            {
                "id": "validate_docx_text",
                "name": "validate_docx_text",
                "kind": "docx_validate_text",
                "inputs": [{"name": "path", "required": True}],
            },
        ],
        "execution_policy": {"timeout_sec": 10},
    }

    result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "operation": "replace_docx_text",
                "arguments": {
                    "input_path": "template.docx",
                    "output_path": "filled.docx",
                    "replacements": {"{{NAME}}": "Ada", "{{ROLE}}": "Engineer"},
                },
            },
            "step_results": {},
        },
    )

    assert result["status"] == "completed"
    assert output.exists()

    validation = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "operation": "validate_docx_text",
                "arguments": {
                    "docx_path": "filled.docx",
                    "must_contain": ["Ada", "Engineer"],
                    "must_not_contain": ["{{NAME}}", "{{ROLE}}"],
                },
            },
            "step_results": {},
        },
    )

    assert validation["status"] == "completed"
    assert validation["final_output"]["ok"] is True


def test_docx_planlet_executes_fill_and_validate(tmp_path: Path) -> None:
    pytest.importorskip("docx")
    from docx import Document

    source = tmp_path / "template.docx"
    output = tmp_path / "filled.docx"
    doc = Document()
    doc.add_paragraph("Hello {{NAME}}")
    doc.save(source)
    spec = {
        "boundary": "typed_operator",
        "operators": [
            {
                "id": "inspect_docx",
                "name": "inspect_docx",
                "kind": "docx_inspect",
                "inputs": [{"name": "path", "required": True, "type": "string"}],
            },
            {
                "id": "replace_docx_text",
                "name": "replace_docx_text",
                "kind": "docx_replace_text",
                "inputs": [
                    {"name": "input_path", "required": True, "type": "string"},
                    {"name": "output_path", "required": True, "type": "string"},
                    {"name": "replacements", "required": True, "type": "object"},
                ],
            },
            {
                "id": "validate_docx_text",
                "name": "validate_docx_text",
                "kind": "docx_validate_text",
                "inputs": [{"name": "path", "required": True, "type": "string"}],
            },
        ],
        "planlets": [
            {
                "id": "docx_template_fill_and_validate",
                "intent": "Fill and validate a DOCX template.",
                "steps": [
                    {
                        "id": "inspect_template",
                        "operation": "inspect_docx",
                        "optional": True,
                        "arguments": {"path": "${input_path}"},
                    },
                    {
                        "id": "fill_template",
                        "operation": "replace_docx_text",
                        "arguments": {
                            "input_path": "${input_path}",
                            "output_path": "${output_path}",
                            "replacements": "${replacements}",
                        },
                    },
                    {
                        "id": "validate_text",
                        "operation": "validate_docx_text",
                        "arguments": {
                            "path": "${output_path}",
                            "must_contain": "${must_contain:[]}",
                            "must_not_contain": "${must_not_contain:[]}",
                        },
                    },
                ],
                "bindings": {
                    "required": ["input_path", "output_path", "replacements"],
                    "optional": ["must_contain", "must_not_contain"],
                    "autobind": {"input_path": ["files[0]"]},
                },
                "keywords": ["docx", "template", "fill"],
            }
        ],
        "execution_policy": {"timeout_sec": 10},
    }

    result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "planlet": "docx_template_fill_and_validate",
                "files": ["template.docx"],
                "arguments": {
                    "output_path": "filled.docx",
                    "replacements": {"{{NAME}}": "Ada"},
                    "must_contain": ["Ada"],
                    "must_not_contain": ["{{NAME}}"],
                },
            },
            "step_results": {},
        },
    )

    assert result["status"] == "completed"
    assert result["runtime_status"] == "planlet_executed"
    assert result["selected_planlet"]["id"] == "docx_template_fill_and_validate"
    assert output.exists()
    assert result["step_results"]["validate_text"]["output"]["ok"] is True


def test_planlet_missing_bindings_returns_guidance(tmp_path: Path) -> None:
    spec = {
        "boundary": "typed_operator",
        "operators": [
            {
                "id": "replace_docx_text",
                "name": "replace_docx_text",
                "kind": "docx_replace_text",
                "inputs": [
                    {"name": "input_path", "required": True, "type": "string"},
                    {"name": "output_path", "required": True, "type": "string"},
                    {"name": "replacements", "required": True, "type": "object"},
                ],
            }
        ],
        "planlets": [
            {
                "id": "docx_template_fill_and_validate",
                "intent": "Fill and validate a DOCX template.",
                "steps": [{"id": "fill", "operation": "replace_docx_text"}],
                "bindings": {"required": ["input_path", "output_path", "replacements"]},
                "keywords": ["docx", "template", "fill"],
            }
        ],
    }

    result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "planlet": "docx_template_fill_and_validate",
                "arguments": {"input_path": "template.docx"},
            },
            "step_results": {},
        },
    )

    assert result["status"] == "completed"
    assert result["runtime_status"] == "planlet_guidance"
    assert result["final_output"]["missing_bindings"] == ["output_path", "replacements"]


def test_operator_argument_aliases_and_type_errors_are_reported(tmp_path: Path) -> None:
    spec = {
        "boundary": "typed_operator",
        "operators": [
            {
                "id": "replace_docx_text",
                "name": "replace_docx_text",
                "kind": "docx_replace_text",
                "inputs": [
                    {"name": "input_path", "required": True, "type": "string"},
                    {"name": "output_path", "required": True, "type": "string"},
                    {"name": "replacements", "required": True, "type": "object"},
                ],
            }
        ],
    }

    result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "operation": "replace_docx_text",
                "arguments": {
                    "path": "template.docx",
                    "output": "filled.docx",
                    "replacements": ["not", "an", "object"],
                },
            },
            "step_results": {},
        },
    )

    assert result["status"] == "blocked"
    details = result["block"]["details"]
    assert details["bound_arguments"]["input_path"] == "template.docx"
    assert details["bound_arguments"]["output_path"] == "filled.docx"
    assert details["errors"] == [
        {"field": "replacements", "reason": "type_mismatch", "expected": "object"}
    ]


def test_pptx_inventory_update_and_validate_generic_operator(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    source = tmp_path / "deck.pptx"
    updated = tmp_path / "updated.pptx"
    inventory = tmp_path / "inventory.json"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.text = "Original title"
    shape_id = shape.shape_id
    prs.save(source)
    spec = {
        "boundary": "typed_operator",
        "operators": [
            {
                "id": "inventory_pptx",
                "name": "inventory_pptx",
                "kind": "pptx_inventory",
                "inputs": [{"name": "input_pptx", "required": True}],
            },
            {
                "id": "update_pptx_text_boxes",
                "name": "update_pptx_text_boxes",
                "kind": "pptx_update_text_boxes",
                "inputs": [
                    {"name": "input_pptx", "required": True},
                    {"name": "output_pptx", "required": True},
                    {"name": "edits", "required": True, "type": "array"},
                ],
            },
            {
                "id": "validate_pptx_file",
                "name": "validate_pptx_file",
                "kind": "pptx_validate_file",
                "inputs": [{"name": "path", "required": True}],
            },
            {
                "id": "add_pptx_reference_slide",
                "name": "add_pptx_reference_slide",
                "kind": "pptx_add_reference_slide",
                "inputs": [
                    {"name": "input_pptx", "required": True},
                    {"name": "output_pptx", "required": True},
                    {"name": "items", "required": True, "type": "array"},
                ],
            },
        ],
        "execution_policy": {"timeout_sec": 10},
    }

    inv_result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "operation": "inventory_pptx",
                "arguments": {"pptx_path": "deck.pptx", "output_json": "inventory.json"},
            },
            "step_results": {},
        },
    )
    assert inv_result["status"] == "completed"
    assert inventory.exists()
    assert inv_result["final_output"]["slides"][0]["shapes"][0]["text"] == "Original title"

    update_result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "operation": "update_pptx_text_boxes",
                "arguments": {
                    "input_path": "deck.pptx",
                    "output_path": "updated.pptx",
                    "targets": [
                        {"slide_index": 1, "shape_id": shape_id, "new_text": "Updated title"}
                    ],
                    "style": {
                        "font_name": "Arial",
                        "font_size": 16,
                        "font_color": "#989596",
                        "bold": False,
                        "alignment": "center",
                        "position": "bottom_center",
                        "fit_to_one_line": True,
                    },
                },
            },
            "step_results": {},
        },
    )
    assert update_result["status"] == "completed"
    assert updated.exists()
    updated_prs = Presentation(str(updated))
    updated_shape = updated_prs.slides[0].shapes[0]
    assert updated_shape.text == "Updated title"
    assert updated_shape.left < Inches(0.5)
    assert updated_shape.width > Inches(9)

    validate_result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "operation": "validate_pptx_file",
                "arguments": {"pptx_path": "updated.pptx"},
            },
            "step_results": {},
        },
    )
    assert validate_result["status"] == "completed"
    assert validate_result["final_output"]["ok"] is True

    referenced = tmp_path / "referenced.pptx"
    add_reference_result = run_boundary_artifact(
        spec,
        {
            "global_inputs": {
                "project_dir": str(tmp_path),
                "operation": "add_pptx_reference_slide",
                "arguments": {
                    "input_path": "updated.pptx",
                    "output_path": "referenced.pptx",
                    "title": "Reference",
                    "titles": ["Paper A", "Paper A", "Paper B"],
                },
            },
            "step_results": {},
        },
    )
    assert add_reference_result["status"] == "completed"
    assert add_reference_result["final_output"]["items_added"] == 2
    assert referenced.exists()
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    with zipfile.ZipFile(referenced) as zipf:
        slide_xml = ET.fromstring(zipf.read("ppt/slides/slide2.xml"))
    assert slide_xml.find(".//a:buAutoNum", ns) is not None
